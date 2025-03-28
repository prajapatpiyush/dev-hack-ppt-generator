from flask import Flask, render_template, request, jsonify, send_file
import os
import google.generativeai as genai  
from pptx import Presentation
import random
import uuid
from datetime import datetime
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)

# Get API key from environment variable
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    raise ValueError(" GEMINI_API_KEY environment variable is not set. Check your .env file.")

# Configure Gemini API with the loaded key
genai.configure(api_key=GEMINI_API_KEY)

def generate_ai_content(prompt):
    """Generates AI-based content using Gemini API"""
    try:
        model = genai.GenerativeModel("gemini-pro")
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f" AI Error: {str(e)}")
        raise Exception(f"Failed to generate AI content: {str(e)}")

def apply_theme(prs, theme):
    """Applies different styles based on theme to all slides"""
    for slide in prs.slides:
        if theme == "Modern":
            # Use RGB tuples for background color
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = (0, 128, 255)
        elif theme == "Classic":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = (255, 215, 0)
        elif theme == "Dark":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = (50, 50, 50)

def parse_ai_content(ai_content):
    """Parses AI content into structured format"""
    content_list = []
    current_heading = None
    current_text = []
    
    for line in ai_content.split("\n"):
        line = line.strip()
        if not line:
            continue
            
        if line.startswith("#") or line.startswith("*"):
            if current_heading and current_text:
                content_list.append({
                    "heading": current_heading,
                    "text": "\n".join(current_text)
                })
            current_heading = line.lstrip("#* ").strip()
            current_text = []
        else:
            current_text.append(line)
    
    if current_heading and current_text:
        content_list.append({
            "heading": current_heading,
            "text": "\n".join(current_text)
        })
    
    return content_list

def create_ppt(title, content_list):
    """Creates a PowerPoint presentation with the given content"""
    try:
        prs = Presentation()
        themes = ["Modern", "Classic", "Dark"]
        selected_theme = random.choice(themes)

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Content slides
        for content in content_list:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            body_placeholder = slide.placeholders[1]
            title_placeholder.text = content["heading"]
            body_placeholder.text = content["text"]

        # Apply theme to all slides
        apply_theme(prs, selected_theme)

        # Generate unique filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]
        ppt_filename = f"presentation_{timestamp}_{unique_id}.pptx"
        
        # Save presentation
        prs.save(ppt_filename)
        return ppt_filename, selected_theme
    except Exception as e:
        print(f" PPT Creation Error: {str(e)}")
        raise Exception(f"Failed to create PowerPoint presentation: {str(e)}")

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    try:
        data = request.get_json()
        if not data or "title" not in data:
            return jsonify({"error": "Title is required"}), 400
            
        title = data.get("title", "AI-Powered Presentation")

        ai_content_prompt = f"""Generate a structured PowerPoint presentation on {title}.
        Format each point with a heading starting with # and detailed explanation.
        Include 3-5 main points with clear headings and explanations."""

        ai_content = generate_ai_content(ai_content_prompt)
        content_list = parse_ai_content(ai_content)

        if not content_list:
            return jsonify({"error": "No content generated"}), 500

        ppt_filename, selected_theme = create_ppt(title, content_list)

        return jsonify({
            "file": ppt_filename,
            "theme": selected_theme,
            "message": " AI-Generated PPT Successfully!"
        })
    
    except Exception as e:
        print(f" Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route("/download_ppt")
def download_ppt():
    try:
        ppt_filename = request.args.get('filename')
        if not ppt_filename:
            return jsonify({"error": "No filename provided"}), 400
        
        return send_file(ppt_filename, as_attachment=True)
    except Exception as e:
        print(f" Download Error: {str(e)}")
        return jsonify({"error": "Failed to download presentation"}), 500

if __name__ == "__main__":
    app.run(debug=True)
